"""Document loader and text splitter for the AI Knowledge Agent."""

import os
from typing import List

from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    TextLoader,
    PyPDFLoader,
    UnstructuredMarkdownLoader,
)

from src.config import DOCS_DIR, CHUNK_SIZE, CHUNK_OVERLAP, SUPPORTED_EXTENSIONS
from src.multimodal import IMAGE_EXTENSIONS

# Map file extensions to their loader class
LOADER_MAP = {
    ".txt": TextLoader,
    ".pdf": PyPDFLoader,
    ".md": UnstructuredMarkdownLoader,
}


def load_docx(filepath: str) -> List[Document]:
    """Load a Word document (.docx) and return as Documents."""
    from docx import Document as DocxDocument

    doc = DocxDocument(filepath)
    text = "\n\n".join(
        paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()
    )

    return [Document(page_content=text, metadata={"source": os.path.basename(filepath), "type": "docx"})]


def load_pptx(filepath: str) -> List[Document]:
    """Load a PowerPoint presentation (.pptx) and return as Documents."""
    from pptx import Presentation

    prs = Presentation(filepath)
    slides = []

    for i, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)

        if texts:
            content = f"[Slide {i}]\n" + "\n".join(texts)
            slides.append(Document(
                page_content=content,
                metadata={"source": os.path.basename(filepath), "type": "pptx", "slide": i},
            ))

    return slides


def load_xlsx(filepath: str) -> List[Document]:
    """Load an Excel spreadsheet (.xlsx) and return as Documents."""
    from openpyxl import load_workbook

    wb = load_workbook(filepath, read_only=True, data_only=True)
    sheets = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip(" |"):
                rows.append(row_text)

        if rows:
            content = f"[Sheet: {sheet_name}]\n" + "\n".join(rows)
            sheets.append(Document(
                page_content=content,
                metadata={"source": os.path.basename(filepath), "type": "xlsx", "sheet": sheet_name},
            ))

    wb.close()
    return sheets


# Custom loaders for Office formats
CUSTOM_LOADERS = {
    ".docx": load_docx,
    ".pptx": load_pptx,
    ".xlsx": load_xlsx,
}


def load_documents(docs_dir: str = DOCS_DIR) -> List[Document]:
    """
    Load all supported documents from the specified directory.
    
    Supports: .txt, .pdf, .md, .docx, .pptx, .xlsx
    
    Args:
        docs_dir: Path to the directory containing documents.
        
    Returns:
        List of Document objects with content and metadata.
    """
    documents = []
    all_supported = SUPPORTED_EXTENSIONS + list(CUSTOM_LOADERS.keys()) + IMAGE_EXTENSIONS

    if not os.path.exists(docs_dir):
        print(f"Documents directory '{docs_dir}' not found. Creating it.")
        os.makedirs(docs_dir)
        return documents

    for filename in os.listdir(docs_dir):
        ext = os.path.splitext(filename)[1].lower()
        if ext not in all_supported:
            print(f"Skipping unsupported file: {filename}")
            continue

        filepath = os.path.join(docs_dir, filename)

        try:
            if ext in CUSTOM_LOADERS:
                loaded = CUSTOM_LOADERS[ext](filepath)
            elif ext in IMAGE_EXTENSIONS:
                from src.multimodal import analyze_image
                description = analyze_image(filepath)
                loaded = [Document(
                    page_content=f"[Image: {filename}]\n\n{description}",
                    metadata={"source": filename, "type": "image", "format": ext.strip(".")},
                )]
            else:
                loader_class = LOADER_MAP.get(ext)
                if loader_class is None:
                    continue
                loader = loader_class(filepath)
                loaded = loader.load()

            for doc in loaded:
                doc.metadata["source"] = filename

            documents.extend(loaded)
            print(f"Loaded: {filename} ({len(loaded)} pages/sections)")
        except Exception as e:
            print(f"Error loading {filename}: {e}")

    print(f"\nTotal documents loaded: {len(documents)}")
    return documents


def split_documents(documents: List[Document]) -> List[Document]:
    """
    Split documents into smaller chunks for vector storage.
    
    Uses RecursiveCharacterTextSplitter which tries to split on
    natural boundaries (paragraphs, sentences) before resorting
    to character-level splits.
    
    Args:
        documents: List of Document objects to split.
        
    Returns:
        List of smaller Document chunks with preserved metadata.
    """
    splitter = RecursiveCharacterTextSplitter(
        chunk_size=CHUNK_SIZE,
        chunk_overlap=CHUNK_OVERLAP,
        length_function=len,
        separators=["\n\n", "\n", ". ", " ", ""],
    )

    chunks = splitter.split_documents(documents)
    print(f"Split into {len(chunks)} chunks")
    return chunks